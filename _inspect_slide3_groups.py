from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

src = r"C:\Users\Aisling Ld Pursuit\Downloads\20260422_Product story template.pptx"
prs = Presentation(src)
slide = prs.slides[2]

def walk(shapes, indent=0):
    for shape in shapes:
        pad = "  " * indent
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{pad}GROUP top={shape.top.inches:.2f} left={shape.left.inches:.2f}")
            walk(shape.shapes, indent + 1)
        elif shape.has_text_frame:
            print(f"{pad}TEXT top={shape.top.inches:.2f} left={shape.left.inches:.2f}: {shape.text!r}")

walk(slide.shapes)
