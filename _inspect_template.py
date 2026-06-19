from pptx import Presentation

src = r"C:\Users\Aisling Ld Pursuit\Downloads\20260422_Product story template.pptx"
prs = Presentation(src)

def fmt(v):
    return f"{v.inches:.2f}" if v is not None else "None"

for si, slide in enumerate(prs.slides, 1):
    print(f"\n=== SLIDE {si} ===")
    shapes = sorted(slide.shapes, key=lambda s: (round((s.top or 0) / 914400, 2), round((s.left or 0) / 914400, 2)))
    for j, shape in enumerate(shapes):
        text = shape.text.replace("\n", " ")[:100] if shape.has_text_frame else ""
        print(
            f"  [{j}] top={fmt(shape.top)} left={fmt(shape.left)} "
            f"w={fmt(shape.width)} h={fmt(shape.height)} type={shape.shape_type}"
        )
        if text:
            print(f"       {text!r}")
