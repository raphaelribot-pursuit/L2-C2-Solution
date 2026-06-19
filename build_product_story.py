"""Rebuild Wisper product story deck from Pursuit template (layout-preserving)."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx.util import Pt

TEMPLATE = Path(r"C:\Users\Aisling Ld Pursuit\Downloads\20260422_Product story template.pptx")
OUT = Path(
    r"C:\Users\Aisling Ld Pursuit\OneDrive\Documents\Pursuit L2 Project\L2 Clone of Wisper\docs\Wisper_Product_Story.pptx"
)
OUT_DOWNLOADS = Path(r"C:\Users\Aisling Ld Pursuit\Downloads\Wisper_Product_Story.pptx")

HEADER_NAME = "Aisling Ld Pursuit & Jimmy"
HEADER_DATE = "June 18, 2026"


def pos(shape) -> tuple[float, float]:
    return (round(shape.left.inches, 2), round(shape.top.inches, 2))


def set_text(shape, text: str, *, font_pt: float | None = None) -> None:
    """Replace placeholder text while keeping run colors, fonts, and box layout."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if p.runs:
        first = p.runs[0]
        first.text = text
        for run in p.runs[1:]:
            run.text = ""
        if font_pt is not None:
            first.font.size = Pt(font_pt)
    else:
        p.text = text
        if font_pt is not None and p.runs:
            p.runs[0].font.size = Pt(font_pt)
    for extra in tf.paragraphs[1:]:
        extra.text = ""


def apply_map(
    slide,
    mapping: dict[tuple[float, float], str],
    *,
    fonts: dict[tuple[float, float], float] | None = None,
) -> None:
    fonts = fonts or {}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        key = pos(shape)
        if key in mapping:
            set_text(shape, mapping[key], font_pt=fonts.get(key))


SLIDE1: dict[tuple[float, float], str] = {
    (0.48, 0.25): HEADER_NAME,
    (7.44, 0.25): HEADER_DATE,
    (0.48, 0.50): (
        "Students lose meeting detail because speech outpaces typing and cloud tools "
        "require audio uploads, resulting in privacy risk and missed notes."
    ),
    (0.48, 2.46): "1. Speech is faster than typing",
    (3.86, 2.46): "2. Cloud tools break trust",
    (7.07, 2.46): "3. Good local options are scarce",
    (0.48, 3.38): (
        "People speak 120–150 words per minute but type ~40 wpm, so live capture "
        "always lags behind the conversation."
    ),
    (3.86, 3.38): (
        "Most transcription apps send audio to the cloud—unacceptable for therapy, "
        "legal, medical, or confidential meetings."
    ),
    (7.07, 3.38): (
        "Whisper Notes proved on-device demand on Mac, but Windows and Linux users "
        "still lack an easy local path."
    ),
}

SLIDE2: dict[tuple[float, float], str] = {
    (0.48, 0.25): HEADER_NAME,
    (7.44, 0.25): HEADER_DATE,
    (0.48, 0.92): "Wisper",
    (0.48, 1.35): (
        "Transcribes audio entirely on your device—no accounts, no cloud upload."
    ),
    (0.48, 3.49): "Capture audio",
    (3.82, 3.49): "Transcribe locally",
    (7.03, 3.49): "Search and export",
    (0.48, 4.20): (
        "Record from mic, import a file, or paste a URL. Welcome guide picks "
        "Small, Medium, or Large for your hardware."
    ),
    (3.82, 4.20): (
        "Whisper.cpp on CPU or GPU (CUDA, Metal, Vulkan). Nothing is uploaded—"
        "transcription stays local."
    ),
    (7.01, 4.20): (
        "SQLite library with full-text search. Export TXT or revisit past "
        "recordings anytime."
    ),
}

SLIDE3: dict[tuple[float, float], str] = {
    (0.48, 0.25): HEADER_NAME,
    (7.44, 0.25): HEADER_DATE,
    (0.48, 0.71): "Under the hood: How Wisper is built",
    (0.56, 1.96): "Mic / file / URL",
    (1.77, 1.96): "GGML speech model",
    (2.98, 1.96): "Language choice",
    (4.19, 1.96): "GPU or CPU",
    (2.25, 3.09): "Wisper",
    (6.12, 2.05): "React UI (Tauri 2)",
    (6.12, 2.75): "wisper-core (Rust + whisper.cpp)",
    (6.12, 3.45): "GPU backends (CUDA / Metal / Vulkan)",
    (6.12, 4.16): "SQLite library + export",
    (0.70, 4.32): "Timestamped transcript",
    (2.25, 4.32): "Searchable library",
    (3.80, 4.32): "TXT export",
}


def verify(slide, mapping: dict[tuple[float, float], str], slide_num: int) -> list[str]:
    issues: list[str] = []
    filled = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        key = pos(shape)
        if key not in mapping:
            continue
        filled += 1
        expected = mapping[key]
        actual = shape.text.strip()
        if actual != expected.strip():
            issues.append(f"Slide {slide_num} @ {key}: got {actual[:60]!r}...")
    if filled != len(mapping):
        issues.append(f"Slide {slide_num}: filled {filled}/{len(mapping)} mapped shapes")
    return issues


def main() -> None:
    from pptx import Presentation

    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT)

    prs = Presentation(str(OUT))
    apply_map(prs.slides[0], SLIDE1, fonts={(0.48, 0.50): 24})
    apply_map(prs.slides[1], SLIDE2)
    apply_map(prs.slides[2], SLIDE3)
    prs.save(str(OUT))

    prs2 = Presentation(str(OUT))
    all_issues: list[str] = []
    for i, (slide, mapping) in enumerate(
        [(prs2.slides[0], SLIDE1), (prs2.slides[1], SLIDE2), (prs2.slides[2], SLIDE3)], 1
    ):
        all_issues.extend(verify(slide, mapping, i))

    downloads_v2 = OUT_DOWNLOADS.with_name("Wisper_Product_Story_v2.pptx")
    for target in (OUT_DOWNLOADS, downloads_v2):
        try:
            shutil.copy2(OUT, target)
            print(f"Copied to {target}")
        except PermissionError:
            print(f"Skipped {target} (file open elsewhere)")

    print(f"Wrote {OUT}")
    if all_issues:
        print("Verification warnings:")
        for issue in all_issues:
            print(" -", issue)
    else:
        print("All mapped placeholders verified.")


if __name__ == "__main__":
    main()
